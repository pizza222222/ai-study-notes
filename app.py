from flask import Flask, request, render_template, jsonify, send_file
import PyPDF2
import google.generativeai as genai
import os
from dotenv import load_dotenv
import io
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
import subprocess
from weasyprint import HTML

# Load environment variables
load_dotenv()

app = Flask(__name__)

# Configure Gemini
genai.configure(api_key=os.environ.get("GEMINI_API_KEY"))
model = genai.GenerativeModel("gemini-2.5-flash")


# Extract text from PDF
def extract_text(file):
    reader = PyPDF2.PdfReader(file)
    text = ""

    for page in reader.pages:
        content = page.extract_text()
        if content:
            text += content

    return text


# Generate structured study notes
def generate_notes(text):

    prompt = f"""
Convert the following study material into well-structured study notes.

Format strictly in HTML using these sections:

<h1>Topic Title</h1>

<h2>Summary</h2>

<h2>Key Concepts</h2>
Use bullet points.

<h2>Definitions</h2>

<h2>Important Questions and Answers</h2>

<h2>Exam Style Questions</h2>

<h2>Book Work / Long Answers</h2>

Use clean academic formatting suitable for students.

At the end write: "Made by SJ"

Study Material:
{text}
"""

    response = model.generate_content(prompt)

    return response.text


# Generate text mind map
def generate_mindmap(text):

    prompt = f"""
Create a Mermaid.js mindmap.

Rules:
- Output ONLY Mermaid code
- Do NOT include explanations
- Start with: mindmap
- Use short keywords

Example:

mindmap
  root((Main Topic))
    Concept1
      Detail1
      Detail2
    Concept2
      Detail3

Topic:
{text}
"""

    response = model.generate_content(prompt)

    diagram = response.text.strip()

    # Remove markdown code blocks if AI adds them
    diagram = diagram.replace("```mermaid", "").replace("```", "").strip()

    return diagram


# Generate PowerPoint slides
def generate_ppt(content):

    prs = Presentation()

    slides = content.split("\n\n")

    for slide_text in slides:

        slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        body = slide.placeholders[1]

        lines = slide_text.split("\n")

        title.text = lines[0]

        if len(lines) > 1:
            body.text = "\n".join(lines[1:])

    file_path = "presentation.pptx"

    prs.save(file_path)

    return file_path


# Convert PPT to PDF
def convert_ppt_to_pdf(ppt_file):

    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        ppt_file
    ])

    return ppt_file.replace(".pptx", ".pdf")


# Home route
@app.route("/", methods=["GET", "POST"])
def index():

    result = ""
    mindmap = ""

    if request.method == "POST":

        if "file" not in request.files:
            return "No file uploaded", 400

        file = request.files["file"]

        if file.filename == "":
            return "No selected file", 400

        text = extract_text(file)

        result = generate_notes(text)

        mindmap = generate_mindmap(text)

    return render_template("index.html", result=result, mindmap=mindmap)



# Download notes as styled PDF
@app.route("/download", methods=["POST"])
def download():

    html_content = request.form.get("content")

    styled_html = f"""
    <html>
    <head>

    <style>

    body {{
        font-family: "Merriweather", serif;
        padding: 40px;
        line-height: 1.8;
        color: #222;
    }}

    h1 {{
        font-size: 34px;
        margin-bottom: 20px;
    }}

    h2 {{
        margin-top: 30px;
        color: #203a43;
    }}

    ul {{
        margin-left: 20px;
    }}

    </style>

    </head>

    <body>

    {html_content}

    </body>

    </html>
    """

    pdf = HTML(string=styled_html).write_pdf()

    return send_file(
        io.BytesIO(pdf),
        download_name="AI_Study_Notes.pdf",
        as_attachment=True,
        mimetype="application/pdf"
    )

# Generate presentation
@app.route("/generate-ppt", methods=["POST"])
def create_ppt():

    content = request.form.get("content")

    ppt_file = generate_ppt(content)

    pdf_file = convert_ppt_to_pdf(ppt_file)

    return send_file(
        pdf_file,
        as_attachment=True,
        download_name="AI_Presentation.pdf"
    )


if __name__ == "__main__":

    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)
